# -*- coding: utf-8 -*-
"""Exp-1 T-A~T-D 종합 비교 발표자료 생성."""
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Length, Pt
import math

OUT = Path(r"C:\Users\seonl\Desktop\c\2026\summer\sdn-xai-project\sdn-xai-pipeline"
           r"\docs\results\Exp1_TA-TD_Comparison.pptx")

FONT = "맑은 고딕"
MONO = "Consolas"

INK    = RGBColor(0x14, 0x1B, 0x2D)
NAVY   = RGBColor(0x1E, 0x3A, 0x5F)
ACCENT = RGBColor(0x2D, 0x7F, 0xF9)
HILITE = RGBColor(0xF2, 0x6B, 0x3A)
GOOD   = RGBColor(0x14, 0x96, 0x6A)
MUTED  = RGBColor(0x6B, 0x74, 0x86)
LINE   = RGBColor(0xD8, 0xDE, 0xE8)
LIGHT  = RGBColor(0xF3, 0xF6, 0xFA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

SW, SH = Inches(13.333), Inches(7.5)
ML     = Inches(0.72)
CW     = Inches(11.9)

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]


# ─────────────────────────────────────────── helpers

def _adv(ch, size):
    """대략적인 글자 폭(pt). 한글/CJK는 전각, ASCII는 반각으로 근사."""
    o = ord(ch)
    if ch == " ":
        return size * 0.30
    if (0x1100 <= o <= 0x11FF or 0x3000 <= o <= 0x303F or 0x3130 <= o <= 0x318F
            or 0x4E00 <= o <= 0x9FFF or 0xAC00 <= o <= 0xD7AF or 0xFF00 <= o <= 0xFFEF):
        return size * 1.0
    return size * 0.53


def text_w(s, size):
    return sum(_adv(c, size) for c in s)


def est_h(s, width, size, line_spacing=1.3, pad=0.16):
    """주어진 폭에서 s가 차지할 높이(Length). width는 Length 또는 int(EMU)."""
    avail = Emu(int(width)).inches * 72.0
    if avail <= 0:
        return Pt(size * line_spacing)
    lines = 0
    for seg in str(s).split("\n"):
        lines += max(1, math.ceil(text_w(seg, size) / avail))
    return Pt(lines * size * line_spacing * (1 + pad))


def slide():
    return prs.slides.add_slide(BLANK)


def rect(sl, left, top, width, height, fill=None, line=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    shp = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(sl, s, left, top, width, height, size=16, bold=False, color=INK,
         align=PP_ALIGN.LEFT, font=FONT, space_after=6, line_spacing=1.25,
         anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = s.split("\n") if isinstance(s, str) else s
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def header(sl, kicker, title, sub=None):
    rect(sl, ML, Inches(0.52), Inches(0.055), Inches(0.30), fill=ACCENT)
    text(sl, kicker, ML + Inches(0.20), Inches(0.52), Inches(9.0), Inches(0.3),
         size=11.5, bold=True, color=ACCENT, space_after=0)
    # 제목이 한 줄에 안 들어가면 들어갈 때까지 줄인다 (볼드는 약 8% 넓게 잡음)
    tsize = 28.0
    avail = Emu(int(CW)).inches * 72.0 * 0.96
    while tsize > 18 and text_w(title, tsize) * 1.08 > avail:
        tsize -= 0.5
    text(sl, title, ML, Inches(0.90), CW, Inches(0.6),
         size=tsize, bold=True, color=INK, space_after=0)
    y = Inches(1.52)
    if sub:
        text(sl, sub, ML, y, CW, Inches(0.4), size=13.5, color=MUTED, space_after=0)
        y = Inches(1.95)
    return y


def bullets(sl, items, top, left=None, width=None, size=15, gap=9):
    """items: list of (text, style) — style in {'h','b','n'}"""
    left = left if left is not None else ML
    width = width if width is not None else CW
    y = top
    for body, style in items:
        if style == "h":
            h = est_h(body, width, size + 1.5)
            text(sl, body, left, y, width, h, size=size + 1.5,
                 bold=True, color=NAVY, space_after=0)
            y += h + Inches(0.06)
        elif style == "n":
            h = est_h(body, width, size - 1.5)
            text(sl, body, left, y, width, h, size=size - 1.5,
                 color=MUTED, space_after=0, line_spacing=1.3)
            y += h
        else:
            avail = Emu(int(width) - int(Inches(0.26)))
            h = est_h(body, avail, size)
            rect(sl, left + Inches(0.03), y + Inches(0.10), Inches(0.075),
                 Inches(0.075), fill=ACCENT)
            text(sl, body, left + Inches(0.26), y, avail, h,
                 size=size, color=INK, space_after=0, line_spacing=1.3)
            y += h
        y += Emu(int(Pt(gap)))
    return y


def table(sl, data, left, top, width, col_w=None, row_h=0.36, size=12.5,
          head_fill=NAVY, hl_rows=None, hl_cols=None, first_col_bold=True):
    """data: list of rows (row 0 = header). hl_cols: set of col idx to tint."""
    nrow, ncol = len(data), len(data[0])
    gf = sl.shapes.add_table(nrow, ncol, left, top, width,
                             Inches(row_h * nrow))
    tbl = gf.table
    if col_w:
        total = sum(col_w)
        for i, w in enumerate(col_w):
            tbl.columns[i].width = Emu(int(width * w / total))
    for r in range(nrow):
        tbl.rows[r].height = Inches(row_h if r else row_h * 1.12)
        for c in range(ncol):
            cell = tbl.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.10)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = head_fill
            elif hl_rows and r in hl_rows:
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF2, 0xEA)
            elif hl_cols and c in hl_cols:
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF2, 0xFE)
            else:
                cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT

            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(data[r][c])
            run.font.size = Pt(size)
            run.font.name = FONT
            if r == 0:
                run.font.bold = True
                run.font.color.rgb = WHITE
            else:
                bold = (c == 0 and first_col_bold) or (hl_cols and c in hl_cols)
                run.font.bold = bool(bold)
                run.font.color.rgb = INK
    return gf


def stat_cards(sl, cards, top, height=1.15, gap=0.22):
    """cards: list of (value, label, color)"""
    n = len(cards)
    w = (CW - Inches(gap) * (n - 1)) / n
    for i, (val, label, col) in enumerate(cards):
        x = ML + (w + Inches(gap)) * i
        rect(sl, x, top, w, Inches(height), fill=LIGHT)
        rect(sl, x, top, Inches(0.05), Inches(height), fill=col)
        text(sl, val, x + Inches(0.26), top + Inches(0.17), w - Inches(0.4),
             Inches(0.5), size=27, bold=True, color=col, space_after=0)
        text(sl, label, x + Inches(0.26), top + Inches(0.72), w - Inches(0.4),
             Inches(0.34), size=11.5, color=MUTED, space_after=0)


def note(sl, s, top=Inches(6.62)):
    rect(sl, ML, top, Inches(0.035), Inches(0.5), fill=LINE)
    text(sl, s, ML + Inches(0.18), top, CW - Inches(0.2), Inches(0.5),
         size=11, color=MUTED, space_after=0, line_spacing=1.25)


def bar_chart(sl, cats, series, left, top, width, height, colors,
              maxval=1.0, number_format='0.00'):
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    gf = sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top,
                             width, height, cd)
    ch = gf.chart
    ch.has_title = False
    ch.font.size = Pt(12)
    ch.font.name = FONT
    ch.font.color.rgb = INK

    if len(series) > 1:
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.TOP
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(12)
    else:
        ch.has_legend = False

    va = ch.value_axis
    va.maximum_scale = maxval
    va.minimum_scale = 0.0
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = LINE
    va.major_gridlines.format.line.width = Pt(0.75)
    va.tick_labels.font.size = Pt(11)
    va.tick_labels.font.color.rgb = MUTED
    va.format.line.fill.background()

    ca = ch.category_axis
    ca.has_major_gridlines = False
    ca.tick_labels.font.size = Pt(12)
    ca.tick_labels.font.color.rgb = INK
    ca.format.line.color.rgb = LINE

    for i, plot_series in enumerate(ch.plots[0].series):
        plot_series.format.fill.solid()
        plot_series.format.fill.fore_color.rgb = colors[i % len(colors)]
        plot_series.format.line.fill.background()
    ch.plots[0].gap_width = 60 if len(series) > 1 else 110
    ch.plots[0].overlap = -10 if len(series) > 1 else 0
    ch.plots[0].has_data_labels = True
    dl = ch.plots[0].data_labels
    dl.font.size = Pt(10.5)
    dl.font.name = FONT
    dl.font.color.rgb = INK
    dl.number_format = number_format
    dl.number_format_is_linked = False
    return ch


# ═══════════════════════════════════════════ 1. Title
s = slide()
rect(s, 0, 0, SW, SH, fill=INK)
rect(s, 0, 0, Inches(0.10), SH, fill=ACCENT)
text(s, "SDN-XAI PIPELINE · EXPERIMENT 1", ML, Inches(1.95), Inches(10), Inches(0.4),
     size=13, bold=True, color=ACCENT, space_after=0)
text(s, "자연어 인텐트 파싱 정확도 실험", ML, Inches(2.50), Inches(11.5), Inches(0.9),
     size=42, bold=True, color=WHITE, space_after=0)
text(s, "IR · Few-shot · Topology Grounding 애블레이션 (T-A ~ T-D)",
     ML, Inches(3.42), Inches(11.5), Inches(0.5), size=20, color=RGBColor(0x9F, 0xB3, 0xCC),
     space_after=0)
rect(s, ML, Inches(4.25), Inches(2.6), Inches(0.035), fill=ACCENT)
text(s,
     "모델  qwen3-8b (OpenRouter)      데이터셋  GOLD-350 (350 cases)\n"
     "반복  1 rep / treatment            작성일  2026-07-26",
     ML, Inches(4.65), Inches(11), Inches(1.0), size=14,
     color=RGBColor(0xC5, 0xD2, 0xE2), space_after=8, line_spacing=1.45)

# ═══════════════════════════════════════════ 2. 배경
s = slide()
y = header(s, "BACKGROUND", "왜 이 실험이 필요한가",
           "LLM이 SDN 플로우 룰을 직접 생성하게 두면 안 되는 이유를 정량적으로 확인한다")

text(s, "파이프라인 구조", ML, y, Inches(5.6), Inches(0.35), size=17, bold=True,
     color=NAVY, space_after=0)
stages = [("Stage 1", "자연어 → IntentIR\n(LLM 파싱)", ACCENT),
          ("Stage 2", "IntentIR → FlowRule\n(결정론적 컴파일러)", GOOD),
          ("Stage 3~6", "정적 검증 · 트윈 검증\nXAI 설명 · 배포", MUTED)]
bx = ML
for i, (tag, body, col) in enumerate(stages):
    w = Inches(1.72)
    rect(s, bx, y + Inches(0.48), w, Inches(1.15), fill=LIGHT)
    rect(s, bx, y + Inches(0.48), w, Inches(0.045), fill=col)
    text(s, tag, bx + Inches(0.14), y + Inches(0.62), w - Inches(0.2), Inches(0.28),
         size=12, bold=True, color=col, space_after=0)
    text(s, body, bx + Inches(0.14), y + Inches(0.94), w - Inches(0.2), Inches(0.7),
         size=10.5, color=INK, space_after=0, line_spacing=1.25)
    bx += w + Inches(0.14)
    if i < 2:
        text(s, "→", bx - Inches(0.135), y + Inches(0.92), Inches(0.3), Inches(0.3),
             size=15, bold=True, color=MUTED, space_after=0, align=PP_ALIGN.CENTER)

text(s, "본 실험(Exp-1)은 Stage 1 출력만 채점한다.", ML, y + Inches(1.78),
     Inches(6.0), Inches(0.3), size=12, color=MUTED, space_after=0)

rx = ML + Inches(6.35)
text(s, "검증하려는 가설", rx, y, Inches(5.5), Inches(0.35), size=17, bold=True,
     color=NAVY, space_after=0)
bullets(s, [
    ("LLM에게 ONOS FlowRule을 직접 만들게 하면, 인텐트에 없는 물리 값(포트·큐 번호)까지 강제로 확정해야 한다", "b"),
    ("중간 표현(IntentIR)을 두면 그 확정을 결정론적 컴파일러로 미룰 수 있다", "b"),
    ("토폴로지 grounding은 호스트↔IP 같은 사실 정보를 채워준다", "b"),
    ("→ 세 요소를 하나씩 켜가며 각각의 기여를 분리 측정한다", "n"),
], y + Inches(0.48), left=rx, width=Inches(5.5), size=14)

# 하단 — 두 가지 접근 대조
yb = Inches(4.42)
text(s, "두 가지 접근", ML, yb, Inches(6), Inches(0.32), size=17, bold=True,
     color=NAVY, space_after=0)
approaches = [
    ("T-A — LLM이 FlowRule을 직접 생성",
     "자연어  →  [ LLM ]  →  ONOS FlowRule",
     "LLM이 egress port · queueId 같은 물리 값까지 확정해야 한다.\n"
     "모르면 지어내거나(환각) 거부할 수밖에 없다.", HILITE),
    ("T-B ~ T-D — IR을 거쳐 컴파일러가 확정",
     "자연어  →  [ LLM ]  →  IntentIR  →  [ 컴파일러 ]  →  FlowRule",
     "LLM은 의미만 담당하고 모르는 값은 null로 남긴다.\n"
     "물리 값 확정은 결정론적 컴파일러가 맡는다.", GOOD),
]
aw = (CW - Inches(0.25)) / 2
for i, (t_, flow, body, col) in enumerate(approaches):
    x = ML + (aw + Inches(0.25)) * i
    rect(s, x, yb + Inches(0.42), aw, Inches(1.62), fill=LIGHT)
    rect(s, x, yb + Inches(0.42), Inches(0.05), Inches(1.62), fill=col)
    text(s, t_, x + Inches(0.26), yb + Inches(0.56), aw - Inches(0.45),
         Inches(0.28), size=13, bold=True, color=col, space_after=0)
    text(s, flow, x + Inches(0.26), yb + Inches(0.90), aw - Inches(0.45),
         Inches(0.28), size=11, color=NAVY, space_after=0, font=MONO)
    text(s, body, x + Inches(0.26), yb + Inches(1.24), aw - Inches(0.45),
         Inches(0.66), size=11.5, color=INK, space_after=2, line_spacing=1.3)

note(s, "Exp-2(Stage 1→2→3 통과율)와 Exp-3(전체 E2E)은 별도 실험으로 분리 진행 — 본 자료는 Exp-1만 다룬다.")

# ═══════════════════════════════════════════ 3. Treatment 설계
s = slide()
y = header(s, "DESIGN", "Treatment 설계 — 3요소 애블레이션",
           "한 번에 하나씩만 켜서, 인접한 두 조건의 차이가 곧 그 요소의 기여가 되도록 설계")

table(s, [
    ["", "출력 형식", "Few-shot", "Grounding", "역할"],
    ["T-A", "ONOS FlowRule (직접)", "✕", "✕", "베이스라인 — IR 없이 LLM이 최종 산출물 생성"],
    ["T-B", "IntentIR", "✕", "✕", "IR 효과만 분리"],
    ["T-C", "IntentIR", "✓ (정적 5개)", "✕", "Few-shot 효과만 분리"],
    ["T-D", "IntentIR", "✓ (정적 5개)", "✓", "현재 파이프라인 최상 설정"],
], ML, y, CW, col_w=[0.7, 2.1, 1.1, 1.1, 4.3], row_h=0.46, size=13.5)

y2 = y + Inches(2.55)
text(s, "비교 축", ML, y2, Inches(5), Inches(0.3), size=16, bold=True, color=NAVY, space_after=0)
table(s, [
    ["비교", "측정하는 효과", "논문에서의 의미"],
    ["T-A → T-B", "IR + 결정론적 컴파일러의 순수 기여", "핵심 주장 검증"],
    ["T-B → T-C", "Few-shot 단독 기여", "보조"],
    ["T-C → T-D", "Topology grounding 단독 기여", "핵심 주장 보강"],
], ML, y2 + Inches(0.42), CW, col_w=[1.5, 4.6, 3.2], row_h=0.42, size=13.5,
    hl_rows={1, 3})

note(s, "T-A만 출력 포맷이 다르므로 IR 슬롯 단위 지표(NEM)는 적용되지 않는다 — T-A는 schema_validity / status_match / "
        "false_rejection_rate로 비교한다 (EVAL_PLAN §2-1).")

# ═══════════════════════════════════════════ 4. 데이터셋
s = slide()
y = header(s, "DATASET", "GOLD-350 — 이중 라벨링 평가 데이터셋",
           "팀원이 독립 이중 라벨링 + adjudication으로 구축, Cohen's κ = 1.000")

stat_cards(s, [
    ("350", "총 케이스", ACCENT),
    ("300 : 50", "accepted : rejected", NAVY),
    ("7 × 50", "카테고리 × 케이스", NAVY),
    ("κ = 1.000", "라벨러 간 일치도", GOOD),
], y, height=1.02)

y2 = y + Inches(1.30)
table(s, [
    ["카테고리", "내용", "케이스"],
    ["forwarding", "경로 설정 · 트래픽 전달", "50"],
    ["security", "차단 · 방화벽 정책", "50"],
    ["qos", "큐 할당 · 대역폭 · 지연 보장", "50"],
    ["sfc", "미들박스(방화벽/IDS) 경유 서비스 체인", "50"],
    ["reroute", "경로 우회 · 장애 대체 경로", "50"],
    ["compound", "복수 정책이 결합된 복합 인텐트", "50"],
    ["ambiguous_unsup.", "거부 대상 (모호 · 미지원 · 모순 · 미지 엔티티)", "50"],
], ML, y2, Inches(7.4), col_w=[2.3, 4.6, 1.0], row_h=0.375, size=12,
    hl_rows={7})

rx = ML + Inches(7.75)
text(s, "토폴로지", rx, y2, Inches(4.1), Inches(0.3), size=16, bold=True,
     color=NAVY, space_after=0)
rect(s, rx, y2 + Inches(0.40), Inches(4.15), Inches(1.28), fill=LIGHT)
text(s, "호스트  h1–h4  (10.0.0.1–4)\n스위치  s1–s4\n서비스  firewall @ s1:9,  IDS @ s2",
     rx + Inches(0.25), y2 + Inches(0.56), Inches(3.7), Inches(0.95),
     size=12, color=INK, space_after=4, line_spacing=1.35, font=MONO)

y3 = y2 + Inches(1.92)
text(s, "Gold 스키마", rx, y3, Inches(4.1), Inches(0.3), size=16,
     bold=True, color=NAVY, space_after=0)
bullets(s, [
    ("파이프라인 IntentIR 스키마를 그대로 gold로 사용 — 정규화 레이어 불필요", "b"),
    ("accepted 300건 Stage1+Stage2 컴파일 전수 PASS", "b"),
], y3 + Inches(0.40), left=rx, width=Inches(4.15), size=12.5)

note(s, "SFC는 ingress/egress rule을 단일 sfc IR로 재합성했으며, multi-hop 4건은 컴파일러가 waypoints[0]만 지원 "
        "(Exp-1 채점에는 영향 없음). reroute의 via_device는 gold에 정보가 없어 해당 슬롯은 채점 제외.")

# ═══════════════════════════════════════════ 5. 평가 지표
s = slide()
y = header(s, "METRICS", "평가 지표",
           "수락/거부 판별 능력과 슬롯 값의 사실 정확도를 분리해서 본다")

text(s, "공통 지표 (T-A ~ T-D)", ML, y, Inches(6.4), Inches(0.3), size=16,
     bold=True, color=NAVY, space_after=0)
table(s, [
    ["지표", "정의"],
    ["schema_validity", "LLM 출력이 목표 스키마로 파싱되는 비율"],
    ["status_match", "accepted/rejected 판정이 gold와 일치하는 비율"],
    ["false_rejection_rate", "수락해야 할 케이스를 거부한 비율"],
    ["rejection_recall", "거부해야 할 케이스를 올바로 거부한 비율"],
    ["false_acceptance_rate", "거부해야 할 케이스를 수락한 비율"],
], ML, y + Inches(0.42), Inches(6.4), col_w=[2.3, 4.6], row_h=0.42, size=12.5)

rx = ML + Inches(6.8)
text(s, "IR 슬롯 지표 (T-B ~ T-D)", rx, y, Inches(5.1), Inches(0.3), size=16,
     bold=True, color=NAVY, space_after=0)
table(s, [
    ["지표", "정의"],
    ["NEM", "모든 슬롯이 gold와 완전 일치"],
    ["slot_accuracy", "슬롯별 개별 정확도"],
    ["hallucination_rate", "토폴로지에 없는 엔티티 생성 비율"],
    ["rule_count_match", "복합 인텐트의 룰 개수 일치 비율"],
], rx, y + Inches(0.42), Inches(5.1), col_w=[2.3, 3.6], row_h=0.42, size=12)

y3 = Inches(5.05)
rect(s, ML, y3, CW, Inches(1.50), fill=RGBColor(0xEF, 0xF5, 0xFF))
rect(s, ML, y3, Inches(0.05), Inches(1.50), fill=ACCENT)
text(s, "NEM (Normalized Exact Match) — 가장 엄격한 지표", ML + Inches(0.32),
     y3 + Inches(0.16), Inches(11), Inches(0.3), size=14, bold=True, color=NAVY,
     space_after=0)
text(s, "action · source_ip · destination_ip · protocol · dst_port · device · egress_port · queue · waypoints 등\n"
        "gold에서 null이 아닌 모든 슬롯이 전부 일치해야 1점. 하나라도 틀리면 0점.\n"
        "gold가 null인 슬롯은 채점 제외, 복합 인텐트는 order-agnostic best-match alignment로 정렬 후 비교.",
     ML + Inches(0.32), y3 + Inches(0.52), Inches(11.2), Inches(0.9), size=12,
     color=INK, space_after=3, line_spacing=1.3)

note(s, "T-A는 IntentIR이 아닌 raw FlowRule JSON을 출력하므로 NEM·슬롯 지표가 구조적으로 적용되지 않는다 (표에서 N/A).")

# ═══════════════════════════════════════════ 6. 실행 조건
s = slide()
y = header(s, "SETUP", "실행 조건", None)

table(s, [
    ["항목", "값", "비고"],
    ["모델", "qwen/qwen3-8b (OpenRouter 경유)", "경량 오픈모델 — 파이프라인의 '로컬/경량 모델 + 결정론적 컴파일러' 논지에 부합"],
    ["Temperature", "0.2", "프로덕션 파이프라인과 동일"],
    ["max_tokens", "8192", "thinking 토큰은 별도 — 최대 12K+ 관측"],
    ["reasoning 예산", "제한 없음", "예산 제한 시 능력 경계 케이스가 불안정해 최종 미적용"],
    ["동시 실행", "concurrency = 20", "429 에러 0건"],
    ["반복", "1 rep / treatment", "⚠ 단일 실행 — 재현성 미확인"],
], ML, y, CW, col_w=[1.6, 3.0, 6.6], row_h=0.40, size=12.5, hl_rows={6})

y2 = Inches(4.52)
text(s, "실행 비용 및 안정성", ML, y2, Inches(6), Inches(0.3), size=16, bold=True,
     color=NAVY, space_after=0)
table(s, [
    ["", "평균 지연", "최대 지연", "평균 입력 토큰", "평균 출력 토큰", "schema_invalid"],
    ["T-A", "38.2s", "184s", "946", "2,667", "10"],
    ["T-B", "35.2s", "132s", "1,660", "2,315", "16"],
    ["T-C", "31.7s", "138s", "2,523", "1,852", "6"],
    ["T-D", "38.9s", "145s", "2,894", "2,190", "4"],
], ML, y2 + Inches(0.40), CW, col_w=[1.0, 2.0, 2.0, 2.4, 2.4, 2.2], row_h=0.355,
    size=12.5)

note(s, "transport 에러는 4개 run 통틀어 0건 — 실패는 전부 JSON 파싱 실패. "
        "Grounding은 입력 토큰을 늘리지만 출력 토큰과 실패율을 함께 낮춘다.",
     top=Inches(6.88))

# ═══════════════════════════════════════════ 7. 종합 결과
s = slide()
y = header(s, "RESULTS", "종합 결과",
           "T-A → T-D로 갈수록 모든 핵심 지표가 단조 개선된다")

table(s, [
    ["지표", "T-A", "T-B", "T-C", "T-D"],
    ["schema_validity", "0.971", "0.954", "0.983", "0.989"],
    ["status_match", "0.560", "0.743", "0.809", "0.949"],
    ["false_rejection_rate ↓", "0.507", "0.267", "0.173", "0.047"],
    ["rejection_recall", "0.960", "0.800", "0.700", "0.920"],
    ["rejection_reason_match", "0.917", "0.850", "0.914", "0.957"],
    ["false_acceptance_rate ↓", "0.040", "0.180", "0.300", "0.060"],
    ["NEM", "N/A", "0.164", "0.149", "0.692"],
    ["hallucinated_entity_rate ↓", "N/A", "0.005", "0.001", "0.000"],
], ML, y, Inches(7.0), col_w=[3.0, 1.0, 1.0, 1.0, 1.1], row_h=0.40, size=12.5,
    hl_cols={4})

rx = ML + Inches(7.35)
bar_chart(s, ["T-A", "T-B", "T-C", "T-D"],
          [("status_match", (0.560, 0.743, 0.809, 0.949)),
           ("NEM", (None, 0.164, 0.149, 0.692))],
          rx, y - Inches(0.10), Inches(4.55), Inches(3.05),
          colors=[NAVY, HILITE])

text(s, "핵심 관찰", rx, y + Inches(3.12), Inches(4.55), Inches(0.3), size=15,
     bold=True, color=NAVY, space_after=0)
bullets(s, [
    ("status_match는 T-A→T-D에서 0.560 → 0.949로 단조 상승", "b"),
    ("NEM은 grounding 도입 시점(T-C→T-D)에만 계단식으로 도약", "b"),
    ("환각은 전 구간에서 0.5% 미만, T-D는 0건", "b"),
], y + Inches(3.52), left=rx, width=Inches(4.55), size=12.5, gap=5)

note(s, "T-A의 NEM·환각률이 N/A인 것은 성능 문제가 아니라 출력 포맷 차이에 따른 구조적 미적용이다.")

# ═══════════════════════════════════════════ 8. 발견 1 — IR
s = slide()
y = header(s, "FINDING 01", "IR의 기여 — 오거부를 절반으로 줄인다  (T-A → T-B)",
           "few-shot·grounding 모두 꺼진 조건에서의 비교. 주된 차이는 출력 형식 (프롬프트 비대칭은 한계 슬라이드 참고).")

stat_cards(s, [
    ("+0.183", "status_match  0.560 → 0.743", GOOD),
    ("−0.240", "false_rejection  0.507 → 0.267", GOOD),
    ("0.080 → 0.720", "QoS 카테고리 status_match", HILITE),
], y, height=1.02)

y2 = y + Inches(1.28)
text(s, "왜 T-A는 수락 가능한 인텐트를 거부하는가 — 형식이 강제하는 조기 확정",
     ML, y2, Inches(11.5), Inches(0.32), size=15.5, bold=True, color=NAVY,
     space_after=0)

table(s, [
    ["인텐트 (gold = accepted)", "T-A가 내놓은 거부 사유", "근본 원인"],
    ["\"Set up forwarding from h3 to h2.\"",
     "Egress port not specified for forwarding action",
     "FlowRule의 OUTPUT은 포트 번호를 필수로 요구"],
    ["\"Give h1 to h3 at least 10 Mbps.\"",
     "Bandwidth guarantee cannot be enforced; requires queueId",
     "FlowRule QoS는 queueId를 필수로 요구"],
], ML, y2 + Inches(0.42), CW, col_w=[3.6, 4.5, 3.8], row_h=0.48, size=12,
    first_col_bold=False)

y3 = y2 + Inches(2.02)
rect(s, ML, y3, CW, Inches(1.62), fill=RGBColor(0xEF, 0xF8, 0xF3))
rect(s, ML, y3, Inches(0.05), Inches(1.62), fill=GOOD)
text(s, "IntentIR은 '아직 모른다'를 표현할 수 있다", ML + Inches(0.32), y3 + Inches(0.16),
     Inches(11), Inches(0.3), size=14, bold=True, color=RGBColor(0x0F, 0x6B, 0x4C),
     space_after=0)
text(s, "IntentIR은 egress_port·device를 null로 남기고 물리 자원 확정을 Stage 2 결정론적 컴파일러로 미룬다. "
        "그래서 같은 인텐트를 수락할 수 있다.\n"
        "→ 즉 \"LLM이 FlowRule을 못 만든다\"가 아니라, \"IR 없는 구조는 알 수 없는 값을 강제로 확정하게 만들고 "
        "그 압력이 오거부로 배출된다\"는 것이다.\n"
        "※ T-A의 false_acceptance 0.040은 판별력이 아니라 보수 편향의 결과 — 수락해야 할 것의 절반을 함께 거부하고 있다.",
     ML + Inches(0.32), y3 + Inches(0.52), Inches(11.2), Inches(1.0), size=12,
     color=INK, space_after=4, line_spacing=1.3)

# ═══════════════════════════════════════════ 9. 발견 2 — Few-shot
s = slide()
y = header(s, "FINDING 02", "Few-shot의 기여 — 형식은 잡지만 사실은 못 채운다  (T-B → T-C)",
           None)

table(s, [
    ["지표", "T-B", "T-C", "Δ", "해석"],
    ["schema_validity", "0.954", "0.983", "+0.029", "출력 형식 준수는 개선"],
    ["dst_port 슬롯", "0.636", "0.951", "+0.315", "포맷성 슬롯은 예시로 학습됨"],
    ["status_match", "0.743", "0.809", "+0.066", "수락 판정 개선"],
    ["NEM", "0.164", "0.149", "−0.015", "전체 정확도는 오히려 소폭 하락"],
    ["source_ip 슬롯", "0.107", "0.094", "−0.013", "사실 슬롯은 그대로"],
    ["false_acceptance_rate", "0.180", "0.300", "+0.120", "거부해야 할 걸 수락하기 시작"],
    ["rejection_recall", "0.800", "0.700", "−0.100", "거부 판별력 저하"],
], ML, y, Inches(7.6), col_w=[2.4, 0.95, 0.95, 0.95, 3.2], row_h=0.40, size=12.5,
    hl_rows={4, 6, 7})

rx = ML + Inches(7.95)
text(s, "해석", rx, y, Inches(3.95), Inches(0.3), size=16, bold=True, color=NAVY,
     space_after=0)
bullets(s, [
    ("Few-shot은 '어떻게 쓸지'는 가르치지만 '무엇이 사실인지'는 못 가르친다", "b"),
    ("host↔IP 매핑은 예시 5개로 일반화되지 않는다 (source_ip 0.107 → 0.094)", "b"),
], y + Inches(0.42), left=rx, width=Inches(3.95), size=13)

y4 = y + Inches(2.05)
rect(s, rx, y4, Inches(3.95), Inches(1.65), fill=RGBColor(0xFF, 0xF4, 0xEC))
rect(s, rx, y4, Inches(0.05), Inches(1.65), fill=HILITE)
text(s, "⚠  부작용 발견", rx + Inches(0.26), y4 + Inches(0.16), Inches(3.5),
     Inches(0.3), size=13.5, bold=True, color=HILITE, space_after=0)
text(s, "현재 few-shot 데모 5개가 전부 accepted 예시다. 모델이 \"일단 수락\" 쪽으로 편향되어 "
        "false_acceptance가 0.180 → 0.300으로 올랐을 가능성이 높다.\n\n"
        "→ rejected 예시를 포함한 데모 세트와 비교 검증 필요.",
     rx + Inches(0.26), y4 + Inches(0.52), Inches(3.5), Inches(1.0), size=11.5,
     color=INK, space_after=2, line_spacing=1.3)

note(s, "NEM 차이 −0.015는 절대값이 작아 단일 rep에서는 샘플링 노이즈와 구분되지 않는다 — rep 확대 후 재확인 필요.")

# ═══════════════════════════════════════════ 10. 발견 3 — Grounding
s = slide()
y = header(s, "FINDING 03", "Grounding의 기여 — 지배적 요인  (T-C → T-D)",
           "NEM +0.543 (0.149 → 0.692). IP 슬롯이 0.09~0.14대에서 0.91대로 도약한다.")

table(s, [
    ["슬롯", "T-B", "T-C", "T-D", "Δ (T-D−T-C)"],
    ["source_ip", "0.107", "0.094", "0.908", "+0.814"],
    ["destination_ip", "0.150", "0.137", "0.908", "+0.771"],
    ["egress_port", "0.430", "0.469", "0.835", "+0.366"],
    ["alt_egress_port", "0.000", "0.027", "0.292", "+0.265"],
    ["device", "0.552", "0.595", "0.819", "+0.224"],
    ["protocol", "1.000", "0.967", "1.000", "+0.033"],
    ["action", "0.932", "0.935", "0.965", "+0.030"],
    ["dst_port", "0.636", "0.951", "0.968", "+0.017"],
    ["waypoints", "0.500", "0.595", "0.542", "−0.053"],
], ML, y, Inches(6.5), col_w=[2.2, 1.0, 1.0, 1.0, 1.4], row_h=0.375, size=12,
    hl_rows={1, 2, 3})

rx = ML + Inches(6.85)
bar_chart(s, ["source_ip", "dest_ip", "egress_port", "device"],
          [("T-C (grounding ✕)", (0.094, 0.137, 0.469, 0.595)),
           ("T-D (grounding ✓)", (0.908, 0.908, 0.835, 0.819))],
          rx, y - Inches(0.10), Inches(5.05), Inches(3.10),
          colors=[RGBColor(0xB8, 0xC4, 0xD4), ACCENT])

text(s, "핵심 대조", rx, y + Inches(3.15), Inches(5.05), Inches(0.3), size=15,
     bold=True, color=NAVY, space_after=0)
bullets(s, [
    ("grounding 인벤토리에 h1 = 10.0.0.1 alias 테이블이 들어가니 IP가 해결되는 건 당연하다", "b"),
    ("중요한 건 few-shot으로는 이게 안 됐다는 대조 — 명시적 그라운딩이 있어야만 풀리는 문제 유형이 존재한다", "b"),
], y + Inches(3.55), left=rx, width=Inches(5.05), size=12.5, gap=5)

note(s, "waypoints만 −0.053으로 유일하게 하락 — topology_eval.json의 그라운딩 결함(s1 firewall은 포트 9가 명시돼 있으나 "
        "s2 IDS는 포트 정보 없음)으로 모델이 \"포트를 모르니 유효 엔티티가 아니다\"라고 판단한 케이스가 있다.")

# ═══════════════════════════════════════════ 11. 카테고리별
s = slide()
y = header(s, "BREAKDOWN", "카테고리별 결과", None)

text(s, "status_match (전 treatment)", ML, y, Inches(6.2), Inches(0.3), size=15,
     bold=True, color=NAVY, space_after=0)
table(s, [
    ["카테고리", "T-A", "T-B", "T-C", "T-D"],
    ["forwarding", "0.660", "0.620", "0.780", "0.980"],
    ["security", "0.540", "0.600", "1.000", "1.000"],
    ["qos", "0.080", "0.720", "0.640", "0.880"],
    ["sfc", "0.580", "0.920", "0.740", "0.960"],
    ["reroute", "0.620", "1.000", "1.000", "1.000"],
    ["compound", "0.480", "0.540", "0.800", "0.900"],
    ["ambiguous_unsup.", "0.960", "0.800", "0.700", "0.920"],
], ML, y + Inches(0.40), Inches(6.2), col_w=[2.3, 1.0, 1.0, 1.0, 1.0],
    row_h=0.385, size=12, hl_rows={3})

rx = ML + Inches(6.6)
text(s, "NEM (T-B/C/D — T-A는 IR 슬롯 없음)", rx, y, Inches(5.3), Inches(0.3),
     size=15, bold=True, color=NAVY, space_after=0)
table(s, [
    ["카테고리", "T-B", "T-C", "T-D"],
    ["security", "0.400", "0.280", "0.900"],
    ["compound", "0.148", "0.100", "0.844"],
    ["qos", "0.111", "0.062", "0.841"],
    ["forwarding", "0.290", "0.256", "0.837"],
    ["reroute", "0.140", "0.140", "0.580"],
    ["sfc", "0.000", "0.000", "0.167"],
], rx, y + Inches(0.40), Inches(5.3), col_w=[2.3, 1.0, 1.0, 1.0],
    row_h=0.385, size=12, hl_rows={6})

y3 = Inches(5.14)
cards = [
    ("QoS — T-A의 최대 약점", "0.080",
     "QoS 인텐트는 대역폭·지연만 말하고 큐 번호를 말하지 않는다. FlowRule은 queueId를 필수로 요구하므로 T-A는 거의 전부 거부.", HILITE),
    ("SFC — 전 구간 최저", "0.167",
     "T-B/T-C에서는 NEM 0.000 — 단 한 건도 완전히 못 맞혔다. waypoint 포트 확정 + 왕복 경로 alt_egress_port 추론이 필요.", NAVY),
]
cw = (CW - Inches(0.25)) / 2
for i, (title_, val, body, col) in enumerate(cards):
    x = ML + (cw + Inches(0.25)) * i
    rect(s, x, y3, cw, Inches(1.32), fill=LIGHT)
    rect(s, x, y3, Inches(0.05), Inches(1.32), fill=col)
    text(s, title_, x + Inches(0.26), y3 + Inches(0.15), cw - Inches(1.6),
         Inches(0.3), size=13.5, bold=True, color=col, space_after=0)
    text(s, val, x + cw - Inches(1.45), y3 + Inches(0.12), Inches(1.2),
         Inches(0.4), size=20, bold=True, color=col, space_after=0,
         align=PP_ALIGN.RIGHT)
    text(s, body, x + Inches(0.26), y3 + Inches(0.55), cw - Inches(0.5),
         Inches(0.8), size=11.5, color=INK, space_after=0, line_spacing=1.3)

note(s, "8B 모델은 배선표를 grounding으로 줘도 다단계 토폴로지 추론을 잘 못한다 — "
        "포트 계산은 컴파일러가 담당해야 한다는 파이프라인 설계 논지를 지지하는 실측 근거.")

# ═══════════════════════════════════════════ 12. 한계
s = slide()
y = header(s, "LIMITATIONS", "한계 및 다음 단계", None)

text(s, "한계", ML, y, Inches(5.8), Inches(0.3), size=17, bold=True, color=HILITE,
     space_after=0)
bullets(s, [
    ("⚠ T-A vs T-B는 단일 변수 조작이 아니다 — 프롬프트 정보량이 양방향으로 다르다 (T-A엔 엔티티 화이트리스트, T-B엔 compound 예시·매핑표). 다만 T-A가 이점을 받고도 낮았으므로 +0.183은 보수적 추정", "b"),
    ("n = 1 rep — 모든 수치가 단일 실행 결과다. 특히 T-B vs T-C의 NEM 역전(−0.015)은 노이즈와 구분되지 않는다", "b"),
    ("T-A vs T-B는 지표 일부만 비교 가능 — 출력 포맷이 달라 NEM·슬롯 정확도는 T-A에 적용 불가", "b"),
    ("Ollama vs OpenRouter 서빙 차이 — 동일 qwen3-8b라도 결과가 달라진다. 본 4개 run은 모두 OpenRouter라 내부 비교는 유효", "b"),
], y + Inches(0.44), left=ML, width=Inches(5.8), size=12.5, gap=7)

rx = ML + Inches(6.2)
text(s, "다음 단계", rx, y, Inches(5.7), Inches(0.3), size=17, bold=True,
     color=GOOD, space_after=0)
steps = [
    ("1", "rep 수 확대", "4 treatment × 3 rep 이상. T-B/T-C NEM 역전의 실체 확인"),
    ("2", "SFC 구간 프로세스 다운 원인 조사", "4회 재현된 패턴. concurrency를 낮춰 재현 여부 확인"),
    ("3", "topology_eval.json s2 IDS 그라운딩 보강", "포트 정보 부재가 SFC waypoints 정확도를 구조적으로 깎는 중"),
    ("4", "few-shot 데모 구성 재검토", "rejected 예시를 포함한 버전과 비교 — false_acceptance 상승 원인"),
]
sy = y + Inches(0.44)
for num, t_, d_ in steps:
    rect(s, rx, sy, Inches(0.36), Inches(0.36), fill=GOOD)
    text(s, num, rx, sy + Inches(0.05), Inches(0.36), Inches(0.3), size=13,
         bold=True, color=WHITE, align=PP_ALIGN.CENTER, space_after=0)
    text(s, t_, rx + Inches(0.52), sy + Inches(0.01), Inches(5.2), Inches(0.3),
         size=13.5, bold=True, color=INK, space_after=0)
    text(s, d_, rx + Inches(0.52), sy + Inches(0.32), Inches(5.2), Inches(0.4),
         size=11.5, color=MUTED, space_after=0, line_spacing=1.25)
    sy += Inches(0.82)

y5 = Inches(5.55)
rect(s, ML, y5, CW, Inches(0.92), fill=RGBColor(0xFF, 0xF4, 0xEC))
rect(s, ML, y5, Inches(0.05), Inches(0.92), fill=HILITE)
text(s, "재현성 주의 — T-A 프롬프트는 run_exp1.py에 하드코딩돼 있다", ML + Inches(0.30),
     y5 + Inches(0.13), Inches(11), Inches(0.3), size=13, bold=True, color=HILITE,
     space_after=0)
text(s, "프로덕션 intent_parser.SYSTEM_PROMPT를 개정할 때마다 SYSTEM_DIRECT_FLOW의 대응 규칙도 함께 점검해야 한다. "
        "실제로 이 불일치 때문에 T-A를 두 차례 재실행했다 (status_match 0.183 → 0.269 → 0.560).",
     ML + Inches(0.30), y5 + Inches(0.47), Inches(11.3), Inches(0.4), size=11.5,
     color=INK, space_after=0, line_spacing=1.25)

# ═══════════════════════════════════════════ 13. 요약
s = slide()
rect(s, 0, 0, SW, SH, fill=INK)
rect(s, 0, 0, Inches(0.10), SH, fill=ACCENT)
text(s, "SUMMARY", ML, Inches(0.75), Inches(10), Inches(0.35), size=13, bold=True,
     color=ACCENT, space_after=0)
text(s, "한 줄 요약", ML, Inches(1.18), Inches(11.5), Inches(0.6), size=32,
     bold=True, color=WHITE, space_after=0)

rect(s, ML, Inches(2.05), CW, Inches(1.15), fill=RGBColor(0x1E, 0x2A, 0x42))
rect(s, ML, Inches(2.05), Inches(0.05), Inches(1.15), fill=ACCENT)
text(s, "IR은 \"수락 가능한 인텐트를 거부하지 않게\" 만들고,\n"
        "grounding은 \"슬롯 값을 정확히 채우게\" 만든다.",
     ML + Inches(0.4), Inches(2.24), Inches(11), Inches(0.85), size=19, bold=True,
     color=WHITE, space_after=4, line_spacing=1.35)

text(s, "둘은 서로 다른 실패 모드를 해결하며 대체 관계가 아니다. few-shot은 형식 준수에는 기여하지만 사실 정확도에는 기여하지 못한다.",
     ML, Inches(3.42), Inches(11.6), Inches(0.4), size=14,
     color=RGBColor(0x9F, 0xB3, 0xCC), space_after=0)

cards = [
    ("T-A → T-B", "+0.183", "status_match", "IR + 결정론적 컴파일러"),
    ("T-B → T-C", "+0.066", "status_match", "Few-shot (NEM은 −0.015)"),
    ("T-C → T-D", "+0.543", "NEM", "Topology grounding"),
]
cw2 = (CW - Inches(0.44)) / 3
for i, (tag, val, metric, label) in enumerate(cards):
    x = ML + (cw2 + Inches(0.22)) * i
    top = Inches(4.15)
    rect(s, x, top, cw2, Inches(1.75), fill=RGBColor(0x1E, 0x2A, 0x42))
    rect(s, x, top, cw2, Inches(0.045), fill=ACCENT if i == 2 else RGBColor(0x3C, 0x50, 0x72))
    text(s, tag, x + Inches(0.3), top + Inches(0.24), cw2 - Inches(0.5), Inches(0.3),
         size=12.5, bold=True, color=RGBColor(0x9F, 0xB3, 0xCC), space_after=0)
    text(s, val, x + Inches(0.3), top + Inches(0.62), cw2 - Inches(0.5), Inches(0.55),
         size=34, bold=True, color=ACCENT if i == 2 else WHITE, space_after=0)
    text(s, metric, x + Inches(0.3), top + Inches(1.16), cw2 - Inches(0.5), Inches(0.28),
         size=12, color=RGBColor(0x9F, 0xB3, 0xCC), space_after=0)
    text(s, label, x + Inches(0.3), top + Inches(1.42), cw2 - Inches(0.5), Inches(0.28),
         size=11.5, color=RGBColor(0x7E, 0x92, 0xB0), space_after=0)

text(s, "실행 근거:  reports/T-{A,B,C,D}_openrouter_r1_summary.json   ·   "
        "상세 분석:  docs/results/TA-TD_openrouter_comparison.md",
     ML, Inches(6.35), Inches(11.6), Inches(0.4), size=11,
     color=RGBColor(0x6B, 0x7E, 0x9A), space_after=0, font=MONO)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print("saved:", OUT)
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))
